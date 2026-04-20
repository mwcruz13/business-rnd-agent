"""
Generate SOC Radar Pattern Library PowerPoint from the HPE Dark template.

Slides:
  1.  Cover — overview stats
  2.  Signal Zones overview (7 zones grid)
  3–9. One slide per signal zone
  10. Disruption Filters (6 filters)
  11. Classification Types (3 types)
  12. Scoring Framework (Impact × Speed)
  13. Action Tiers
  14. RPV Framework
  15. Evidence Capture + Confidence Ratings

Usage:
  cat generate_soc_radar_pptx.py | docker compose exec -T bmi-backend python -
"""
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# --- HPE Dark Theme Colors ---
BG_DARK     = RGBColor(0x29, 0x2D, 0x3A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
MUTED_TEXT  = RGBColor(0xB1, 0xB9, 0xBE)
MID_GRAY    = RGBColor(0x53, 0x5C, 0x66)
HPE_GREEN   = RGBColor(0x00, 0xE0, 0xAF)
HPE_PURPLE  = RGBColor(0x9B, 0x84, 0xFC)
HPE_BLUE    = RGBColor(0x65, 0xAE, 0xF9)
HPE_CYAN    = RGBColor(0x62, 0xE5, 0xF6)
HPE_ORANGE  = RGBColor(0xF2, 0x6B, 0x43)
HPE_YELLOW  = RGBColor(0xF2, 0xC9, 0x4C)
CARD_BG     = RGBColor(0x33, 0x3A, 0x4A)

FONT_HEADING = "HPE Graphik Semibold"
FONT_BODY    = "HPE Graphik"

# Zone accent colors
ZONE_COLORS = [HPE_GREEN, HPE_ORANGE, HPE_BLUE, HPE_PURPLE, HPE_CYAN, HPE_YELLOW,
               RGBColor(0xFF, 0x6B, 0x9D)]

# --- Paths ---
TEMPLATE   = "hpe_dark_template.pptx"
DATA_FILE  = "backend/app/patterns/soc-radar-pattern-library.json"
OUTPUT     = "reports/SOC_Radar_Pattern_Library.pptx"

# ---- LOAD DATA ----
with open(DATA_FILE, encoding="utf-8") as f:
    data = json.load(f)

# ---- INIT PRESENTATION ----
prs = Presentation(TEMPLATE)

sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst):
    rId = sldId.get(qn("r:id"))
    if rId:
        try:
            prs.part.drop_rel(rId)
        except KeyError:
            pass
    sldIdLst.remove(sldId)

try:
    BLANK_LAYOUT = prs.slide_layouts[37]
except IndexError:
    BLANK_LAYOUT = prs.slide_layouts[-1]


# ---- HELPERS ----
def add_slide():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG_DARK
    return slide


def text_box(slide, left, top, width, height, text, *,
             font=FONT_BODY, size=14, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox


def rounded_rect(slide, left, top, width, height, fill, *, border=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


# ============================================================
# SLIDE 1 — COVER
# ============================================================
s = add_slide()
rounded_rect(s, 4.0, 0.8, 5.3, 0.45, HPE_ORANGE)
text_box(s, 4.0, 0.82, 5.3, 0.45, "SIGNALS OF CHANGE",
         font=FONT_HEADING, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

text_box(s, 1.0, 1.8, 11.3, 0.8, "SOC Radar",
         font=FONT_HEADING, size=40, color=HPE_ORANGE, bold=True, align=PP_ALIGN.CENTER)
text_box(s, 1.0, 2.5, 11.3, 0.8, "Disruption Pattern Library",
         font=FONT_HEADING, size=32, color=HPE_BLUE, bold=True, align=PP_ALIGN.CENTER)

text_box(s, 2.5, 3.6, 8.3, 1.0,
         "Detect, classify, and score signals of change using the canonical "
         "frameworks from Seeing What's Next — 7 signal zones, 6 disruption "
         "filters, and the RPV assessment framework.",
         size=13, color=MUTED_TEXT, align=PP_ALIGN.CENTER)

stats = [
    ("7", "Signal Zones", HPE_GREEN),
    ("6", "Disruption Filters", HPE_ORANGE),
    ("3", "Classification Types", HPE_PURPLE),
    ("3", "RPV Dimensions", HPE_CYAN),
]
for i, (num, label, clr) in enumerate(stats):
    x = 2.0 + i * 2.6
    text_box(s, x, 5.0, 2.0, 0.6, num, font=FONT_HEADING, size=36,
             color=clr, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, x, 5.6, 2.0, 0.3, label, size=9, color=MUTED_TEXT,
             align=PP_ALIGN.CENTER)

text_box(s, 0.5, 6.8, 12.3, 0.4,
         "Source: Scott D. Anthony, Clayton Christensen, Erik Roth — Seeing What's Next",
         size=9, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2 — SIGNAL ZONES OVERVIEW
# ============================================================
zones = data["signal_zones"]["zones"]
s = add_slide()
text_box(s, 0.6, 0.3, 5, 0.3, "SIGNAL ZONES",
         font=FONT_HEADING, size=10, color=HPE_GREEN, bold=True)
text_box(s, 0.6, 0.6, 10, 0.5, "7 Zones for Detecting Signals of Change",
         font=FONT_HEADING, size=22, color=WHITE, bold=True)

# Grid: 3 + 3 + 1
card_w = 3.9
card_h = 1.5
gap = 0.2
start_x = 0.6
start_y = 1.3

for idx, zone in enumerate(zones):
    col = idx % 3
    row = idx // 3
    cx = start_x + col * (card_w + gap)
    cy = start_y + row * (card_h + gap)
    zcolor = ZONE_COLORS[idx]

    rounded_rect(s, cx, cy, card_w, card_h, CARD_BG, border=MID_GRAY)
    # Accent bar
    shape = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(cx), Inches(cy), Inches(card_w), Inches(0.04),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = zcolor
    shape.line.fill.background()

    text_box(s, cx + 0.12, cy + 0.12, card_w - 0.24, 0.25,
             zone["name"], font=FONT_HEADING, size=11, color=zcolor, bold=True)
    text_box(s, cx + 0.12, cy + 0.4, card_w - 0.24, 0.55,
             zone["description"], size=8, color=MUTED_TEXT)
    text_box(s, cx + 0.12, cy + 0.95, card_w - 0.24, 0.15,
             "WATCH FOR", size=7, color=MID_GRAY, bold=True)
    text_box(s, cx + 0.12, cy + 1.1, card_w - 0.24, 0.35,
             zone["watch_for"], size=7, color=MUTED_TEXT)


# ============================================================
# SLIDES 3–9 — ONE PER ZONE
# ============================================================
for idx, zone in enumerate(zones):
    s = add_slide()
    zcolor = ZONE_COLORS[idx]

    text_box(s, 0.6, 0.3, 5, 0.3, f"SIGNAL ZONE {idx + 1} OF 7",
             font=FONT_HEADING, size=10, color=zcolor, bold=True)
    text_box(s, 0.6, 0.7, 10, 0.6, zone["name"],
             font=FONT_HEADING, size=30, color=zcolor, bold=True)

    # Description
    text_box(s, 0.6, 1.6, 11.5, 0.6, zone["description"],
             size=16, color=MUTED_TEXT)

    # Watch for section
    rounded_rect(s, 0.6, 2.6, 11.5, 2.5, CARD_BG)
    shape = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(2.6), Inches(0.06), Inches(2.5),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = zcolor
    shape.line.fill.background()

    text_box(s, 0.9, 2.75, 3.0, 0.3, "WATCH FOR",
             font=FONT_HEADING, size=10, color=zcolor, bold=True)
    text_box(s, 0.9, 3.15, 10.8, 1.8, zone["watch_for"],
             size=14, color=MUTED_TEXT)


# ============================================================
# SLIDE 10 — DISRUPTION FILTERS
# ============================================================
filters = data["disruption_filters"]["filters"]
s = add_slide()
text_box(s, 0.6, 0.3, 5, 0.3, "DISRUPTION FILTERS",
         font=FONT_HEADING, size=10, color=HPE_ORANGE, bold=True)
text_box(s, 0.6, 0.6, 10, 0.5, "6 Filters for Assessing Disruptive Potential",
         font=FONT_HEADING, size=22, color=WHITE, bold=True)

# Threshold note
rounded_rect(s, 0.6, 1.2, 12.1, 0.4, RGBColor(0x3D, 0x2E, 0x28), border=HPE_ORANGE)
text_box(s, 0.8, 1.22, 11.7, 0.36, data["disruption_filters"]["threshold_note"],
         size=10, color=HPE_ORANGE)

y = 1.8
for i, f in enumerate(filters):
    rounded_rect(s, 0.6, y, 12.1, 0.7, CARD_BG, border=MID_GRAY)
    # Orange left bar
    shape = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(y), Inches(0.06), Inches(0.7),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = HPE_ORANGE
    shape.line.fill.background()

    text_box(s, 0.85, y + 0.05, 3.5, 0.25,
             f"{i+1}. {f['name']}", font=FONT_HEADING, size=11,
             color=HPE_ORANGE, bold=True)
    text_box(s, 0.85, y + 0.32, 11.5, 0.35,
             f['question'], size=9, color=MUTED_TEXT)
    y += 0.78


# ============================================================
# SLIDE 11 — CLASSIFICATION TYPES
# ============================================================
types = data["classification_types"]
s = add_slide()
text_box(s, 0.6, 0.3, 5, 0.3, "CLASSIFICATION",
         font=FONT_HEADING, size=10, color=HPE_PURPLE, bold=True)
text_box(s, 0.6, 0.6, 10, 0.5, "3 Signal Classification Types",
         font=FONT_HEADING, size=22, color=WHITE, bold=True)

type_colors = [HPE_GREEN, HPE_ORANGE, HPE_PURPLE]
y = 1.3
for i, t in enumerate(types):
    tc = type_colors[i]
    rounded_rect(s, 0.6, y, 12.1, 1.4, CARD_BG, border=MID_GRAY)
    # Top accent bar
    shape = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(y), Inches(12.1), Inches(0.04),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = tc
    shape.line.fill.background()

    text_box(s, 0.85, y + 0.12, 5, 0.3, t["name"],
             font=FONT_HEADING, size=16, color=tc, bold=True)
    text_box(s, 0.85, y + 0.5, 11.5, 0.4, t["description"],
             size=11, color=MUTED_TEXT)
    text_box(s, 0.85, y + 0.9, 2.0, 0.2, "STRATEGIC RESPONSE",
             size=8, color=MID_GRAY, bold=True)
    text_box(s, 3.0, y + 0.9, 9.5, 0.4, t["strategic_response"],
             size=11, color=HPE_CYAN)
    y += 1.55


# ============================================================
# SLIDE 12 — SCORING FRAMEWORK
# ============================================================
scoring = data["scoring"]
s = add_slide()
text_box(s, 0.6, 0.3, 5, 0.3, "SCORING FRAMEWORK",
         font=FONT_HEADING, size=10, color=HPE_BLUE, bold=True)
text_box(s, 0.6, 0.6, 10, 0.5, "Impact × Speed Priority Matrix",
         font=FONT_HEADING, size=22, color=WHITE, bold=True)

# Formula box
rounded_rect(s, 3.5, 1.3, 6.3, 0.6, CARD_BG, border=HPE_ORANGE)
text_box(s, 3.5, 1.32, 6.3, 0.56, scoring["formula"],
         font=FONT_HEADING, size=20, color=HPE_ORANGE, bold=True, align=PP_ALIGN.CENTER)

# Dimensions
y = 2.2
for dim in scoring["dimensions"]:
    text_box(s, 0.6, y, 3, 0.3, f"{dim['name']} ({dim['scale']})",
             font=FONT_HEADING, size=14, color=HPE_BLUE, bold=True)
    y += 0.35
    for lvl, desc in dim["levels"].items():
        rounded_rect(s, 0.8, y, 11.7, 0.35, CARD_BG)
        text_box(s, 1.0, y + 0.03, 0.5, 0.3, lvl,
                 font=FONT_HEADING, size=14, color=HPE_BLUE, bold=True,
                 align=PP_ALIGN.CENTER)
        text_box(s, 1.6, y + 0.04, 10.7, 0.28, desc,
                 size=10, color=MUTED_TEXT)
        y += 0.4
    y += 0.15


# ============================================================
# SLIDE 13 — ACTION TIERS
# ============================================================
s = add_slide()
text_box(s, 0.6, 0.3, 5, 0.3, "ACTION TIERS",
         font=FONT_HEADING, size=10, color=HPE_ORANGE, bold=True)
text_box(s, 0.6, 0.6, 10, 0.5, "Priority Score → Response Action",
         font=FONT_HEADING, size=22, color=WHITE, bold=True)

tier_colors = [HPE_BLUE, HPE_YELLOW, HPE_ORANGE]
tiers = scoring["action_tiers"]
card_w = 3.8
gap = 0.35
start_x = 0.8

for i, tier in enumerate(tiers):
    tc = tier_colors[i]
    cx = start_x + i * (card_w + gap)
    rounded_rect(s, cx, 1.5, card_w, 3.5, CARD_BG, border=MID_GRAY)

    # Top accent bar
    shape = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(cx), Inches(1.5), Inches(card_w), Inches(0.06),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = tc
    shape.line.fill.background()

    text_box(s, cx, 1.75, card_w, 0.5, tier["score_range"],
             font=FONT_HEADING, size=36, color=tc, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, cx, 2.3, card_w, 0.4, tier["tier"],
             font=FONT_HEADING, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, cx + 0.2, 2.85, card_w - 0.4, 0.8, tier["action"],
             size=12, color=MUTED_TEXT, align=PP_ALIGN.CENTER)
    text_box(s, cx, 3.8, card_w, 0.3, tier["timeframe"],
             size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 14 — RPV FRAMEWORK
# ============================================================
rpv = data["rpv_framework"]
s = add_slide()
text_box(s, 0.6, 0.3, 5, 0.3, "RPV FRAMEWORK",
         font=FONT_HEADING, size=10, color=HPE_CYAN, bold=True)
text_box(s, 0.6, 0.6, 10, 0.5, rpv["name"],
         font=FONT_HEADING, size=22, color=WHITE, bold=True)
text_box(s, 0.6, 1.1, 11, 0.4, rpv["purpose"],
         size=12, color=MUTED_TEXT)

rpv_colors = [HPE_GREEN, HPE_BLUE, HPE_PURPLE]
rpv_icons = ["R", "P", "V"]
card_w = 3.8
gap = 0.35
start_x = 0.8

for i, dim in enumerate(rpv["dimensions"]):
    rc = rpv_colors[i]
    cx = start_x + i * (card_w + gap)
    rounded_rect(s, cx, 1.8, card_w, 3.0, CARD_BG, border=MID_GRAY)

    # Icon circle
    text_box(s, cx, 2.0, card_w, 0.5, rpv_icons[i],
             font=FONT_HEADING, size=36, color=rc, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, cx, 2.6, card_w, 0.3, dim["name"],
             font=FONT_HEADING, size=16, color=rc, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, cx + 0.2, 3.1, card_w - 0.4, 1.5, dim["question"],
             size=11, color=MUTED_TEXT, align=PP_ALIGN.CENTER)

# Interpretation
rounded_rect(s, 0.6, 5.2, 12.1, 0.8, RGBColor(0x2A, 0x34, 0x3D), border=HPE_CYAN)
text_box(s, 0.8, 5.25, 11.7, 0.7, rpv["interpretation"],
         size=11, color=HPE_CYAN)


# ============================================================
# SLIDE 15 — EVIDENCE CAPTURE + CONFIDENCE
# ============================================================
s = add_slide()
text_box(s, 0.6, 0.3, 5, 0.3, "EVIDENCE & CONFIDENCE",
         font=FONT_HEADING, size=10, color=HPE_CYAN, bold=True)
text_box(s, 0.6, 0.6, 10, 0.5, "Evidence Capture Fields & Confidence Ratings",
         font=FONT_HEADING, size=22, color=WHITE, bold=True)

# Evidence fields
fields = data["evidence_capture"]["required_fields"]
text_box(s, 0.6, 1.2, 5, 0.25, "REQUIRED EVIDENCE FIELDS",
         font=FONT_HEADING, size=10, color=HPE_CYAN, bold=True)

col_w = 3.8
for i, field in enumerate(fields):
    col = i % 3
    row = i // 3
    cx = 0.6 + col * (col_w + 0.2)
    cy = 1.55 + row * 0.45
    rounded_rect(s, cx, cy, col_w, 0.38, CARD_BG, border=MID_GRAY)
    text_box(s, cx + 0.1, cy + 0.04, col_w - 0.2, 0.3, field,
             size=11, color=HPE_CYAN, bold=True, align=PP_ALIGN.CENTER)

# Confidence ratings
conf = data["confidence_ratings"]
text_box(s, 0.6, 2.8, 5, 0.25, "CONFIDENCE RATINGS",
         font=FONT_HEADING, size=10, color=HPE_GREEN, bold=True)

conf_colors = [HPE_GREEN, HPE_YELLOW, HPE_ORANGE]
for i, c in enumerate(conf):
    cc = conf_colors[i]
    cy = 3.15 + i * 1.1
    rounded_rect(s, 0.6, cy, 12.1, 0.95, CARD_BG, border=MID_GRAY)
    text_box(s, 0.85, cy + 0.05, 1.5, 0.3, c["level"],
             font=FONT_HEADING, size=16, color=cc, bold=True)
    text_box(s, 2.5, cy + 0.05, 4.0, 0.3, c["meaning"],
             size=11, color=WHITE)
    text_box(s, 0.85, cy + 0.45, 11.5, 0.45, c["guidance"],
             size=9, color=MUTED_TEXT)


# ---- SAVE ----
Path("reports").mkdir(exist_ok=True)
prs.save(OUTPUT)
print(f"Saved {OUTPUT} with {len(prs.slides)} slides")
