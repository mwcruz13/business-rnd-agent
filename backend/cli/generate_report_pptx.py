"""
Generate a BMI Workflow Report PowerPoint from a completed workflow session.

Reads session state from the database and produces a presentation using
the HPE Dark Inspirational R9 template.

Usage (inside container):
    python -m backend.cli.generate_report_pptx --session-id <ID>
    python -m backend.cli.generate_report_pptx   # uses most recent completed session

Output: reports/<session_id>_report.pptx  (or --output <path>)
"""
from __future__ import annotations

import io
import json
import re
import textwrap
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── HPE Dark Theme ──────────────────────────────────────────────────────
BG_DARK     = RGBColor(0x29, 0x2D, 0x3A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
MUTED       = RGBColor(0xB1, 0xB9, 0xBE)
MID_GRAY    = RGBColor(0x53, 0x5C, 0x66)
HPE_GREEN   = RGBColor(0x00, 0xE0, 0xAF)
HPE_GREEN_D = RGBColor(0x00, 0xB8, 0x8C)
HPE_PURPLE  = RGBColor(0x9B, 0x84, 0xFC)
HPE_BLUE    = RGBColor(0x65, 0xAE, 0xF9)
HPE_CYAN    = RGBColor(0x62, 0xE5, 0xF6)
HPE_ORANGE  = RGBColor(0xF2, 0x6B, 0x43)
HPE_YELLOW  = RGBColor(0xFF, 0xBC, 0x44)
CARD_BG     = RGBColor(0x33, 0x3A, 0x4A)
CARD_BG2    = RGBColor(0x2F, 0x35, 0x44)
CARD_BORDER = MID_GRAY
FADED_NUM   = RGBColor(0x44, 0x4C, 0x5A)

FH = "HPE Graphik Semibold"
FB = "HPE Graphik"

TEMPLATE = "backend/assets/hpe_dark_template.pptx"
ILLUSTRATIONS_PPTX = "backend/assets/hpe_spot_illustrations.pptx"
BLANK_LAYOUT = 37

TIER_COLORS = {"Act": HPE_GREEN, "Investigate": HPE_BLUE, "Monitor": HPE_YELLOW, "Ignore": MID_GRAY}
CAT_COLORS  = {"Desirability": HPE_GREEN, "Feasibility": HPE_BLUE, "Viability": HPE_PURPLE}

# ── Illustration Mapping ────────────────────────────────────────────────
SLIDE_ILLUSTRATIONS: dict[str, list[str]] = {
    "cover":        ["AI", "Enterprise"],
    "signals":      ["Observability", "Data"],
    "priority":     ["Performance", "Scalability"],
    "pattern":      ["Market Growth", "Scalability"],
    "profile":      ["Partnership", "Enterprise"],
    "vdt":          ["Cost Savings", "Time Savings"],
    "insights":     ["Data", "Agentic AI"],
    "insights2":    ["Services", "Partnership"],
    "vpc":          ["Services", "Cloud"],
    "bmc":          ["Financial Services", "Infrastructure"],
    "fit":          ["Orchestration", "Data Governance"],
    "assumptions":  ["Observability", "Security"],
    "dvf":          ["Networking", "Orchestration"],
    "experiments":  ["Quick Setup", "Machine Learning"],
    "worksheets":   ["Machine Learning", "Quick Setup"],
}


# ── Helpers ─────────────────────────────────────────────────────────────

def _init_prs(template_path: str) -> Presentation:
    prs = Presentation(template_path)
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        if rId:
            try:
                prs.part.drop_rel(rId)
            except KeyError:
                pass
        sldIdLst.remove(sldId)
    return prs


def _slide(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG_DARK
    return s


def _tb(s, l, t, w, h, text, font=FB, sz=12, clr=WHITE, bold=False,
        align=PP_ALIGN.LEFT, spacing_after=0):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing_after:
        p.space_after = Pt(spacing_after)
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(sz)
    r.font.color.rgb = clr
    r.font.bold = bold
    return box


def _multi_tb(s, l, t, w, h, lines, font=FB, sz=11, clr=MUTED, bold=False,
              align=PP_ALIGN.LEFT, sp=3, bullet=False):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sp)
        r = p.add_run()
        r.text = ("\u2022  " + line) if bullet else line
        r.font.name = font
        r.font.size = Pt(sz)
        r.font.color.rgb = clr
        r.font.bold = bold
    return box


def _rect(s, l, t, w, h, fill, border=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = s.shapes.add_shape(shape_type, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh


def _bar(s, l, t, w, h, fill):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


def _footer(s, text="BMI WORKFLOW REPORT"):
    _tb(s, 0.5, 6.9, 12.3, 0.35, text, FB, 10, MID_GRAY, False, PP_ALIGN.CENTER)


def _section_label(s, y, text):
    _tb(s, 0.7, y, 8, 0.3, text, FH, 12, HPE_GREEN, True)


def _section_title(s, y, text, sz=24):
    _tb(s, 0.7, y + 0.28, 11, 0.55, text, FH, sz, WHITE, True)


def _notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def _clean(text: str) -> str:
    """Strip markdown artifacts but keep full content."""
    if not text:
        return ""
    text = text.strip()
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    return text


def _parse_md_table(md_text: str) -> list[dict[str, str]]:
    """Parse a markdown table into a list of dicts. Returns [] if no table found."""
    lines = [l.strip() for l in md_text.strip().split("\n") if l.strip()]
    table_lines = [l for l in lines if l.startswith("|")]
    if len(table_lines) < 3:
        return []
    headers_raw = [h.strip() for h in table_lines[0].split("|") if h.strip()]
    rows = []
    for row_line in table_lines[2:]:
        cells = [c.strip() for c in row_line.split("|") if c.strip() != ""]
        if cells:
            row = {}
            for i, h in enumerate(headers_raw):
                row[h] = cells[i] if i < len(cells) else ""
            rows.append(row)
    return rows


# ── Table Helper ────────────────────────────────────────────────────────

def _styled_table(s, left, top, width, col_widths, headers, rows,
                  accent=HPE_GREEN, font_sz=11, header_sz=11):
    """Add an HPE-themed table to a slide.

    Args:
        s: slide object
        left, top, width: position and total width in inches
        col_widths: list of proportional widths (will be normalized to sum=width)
        headers: list of header strings
        rows: list of lists of cell strings
        accent: header row background color
        font_sz: body font size in pt
        header_sz: header font size in pt
    Returns:
        The table shape
    """
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)
    row_height = Inches(0.35)

    tbl_shape = s.shapes.add_table(n_rows, n_cols,
                                    Inches(left), Inches(top),
                                    Inches(width), row_height * n_rows)
    tbl = tbl_shape.table

    # Normalize column widths
    total = sum(col_widths)
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = Inches(width * cw / total)

    # Style header row
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = hdr
        r.font.name = FH
        r.font.size = Pt(header_sz)
        r.font.color.rgb = WHITE
        r.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        # Remove borders
        _cell_borders(cell, MID_GRAY)

    # Style data rows
    for ri, row_data in enumerate(rows):
        bg = CARD_BG if ri % 2 == 0 else CARD_BG2
        for ci in range(n_cols):
            cell = tbl.cell(ri + 1, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = row_data[ci] if ci < len(row_data) else ""
            r.font.name = FB
            r.font.size = Pt(font_sz)
            r.font.color.rgb = MUTED
            p.alignment = PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.vertical_anchor = MSO_ANCHOR.TOP
            cell.text_frame.word_wrap = True
            _cell_borders(cell, MID_GRAY)

    return tbl_shape


def _cell_borders(cell, color):
    """Set thin borders on a table cell using low-level XML."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ("lnL", "lnR", "lnT", "lnB"):
        ln = tcPr.find(qn(f"a:{edge}"))
        if ln is None:
            from lxml import etree
            ln = etree.SubElement(tcPr, qn(f"a:{edge}"), attrib={"w": "6350", "cap": "flat", "cmpd": "sng"})
        else:
            ln.set("w", "6350")
        # Add or replace solidFill
        sf = ln.find(qn("a:solidFill"))
        if sf is None:
            from lxml import etree
            sf = etree.SubElement(ln, qn("a:solidFill"))
        else:
            sf.clear()
        from lxml import etree
        srgb = etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": f"{color}"})


# ── Illustration Helpers ────────────────────────────────────────────────

def _load_illustrations(pptx_path: str) -> dict[str, bytes]:
    """Extract spot illustrations from the HPE illustrations deck."""
    if not Path(pptx_path).exists():
        return {}
    try:
        prs = Presentation(pptx_path)
    except Exception:
        return {}

    illustrations: dict[str, bytes] = {}
    for si in [4, 5, 6]:
        if si >= len(prs.slides):
            continue
        slide = prs.slides[si]
        images = []
        labels = []
        for shape in slide.shapes:
            try:
                img = shape.image
                images.append({"blob": img.blob, "left": shape.left, "top": shape.top})
            except Exception:
                pass
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                skip_prefixes = ("Spot illustrations", "Confidential")
                if t and not t.isdigit() and not any(t.startswith(s) for s in skip_prefixes):
                    labels.append({"text": t, "left": shape.left, "top": shape.top})
        images.sort(key=lambda x: (x["top"], x["left"]))
        labels.sort(key=lambda x: (x["top"], x["left"]))
        for idx, img in enumerate(images):
            if idx < len(labels):
                illustrations[labels[idx]["text"]] = img["blob"]
    return illustrations


def _place_illustration(s, illustrations: dict[str, bytes],
                        labels: list[str],
                        left: float = 11.8, top: float = 0.25, size: float = 0.8):
    """Place one or two spot illustrations on a slide."""
    if not illustrations:
        return
    offset = 0.0
    for label in labels:
        blob = illustrations.get(label)
        if blob:
            stream = io.BytesIO(blob)
            s.shapes.add_picture(stream, Inches(left), Inches(top + offset),
                                 Inches(size), Inches(size))
            offset += size + 0.1


# ── Slide Builders ──────────────────────────────────────────────────────

def _slide_cover(prs, state, meta, illus):
    s = _slide(prs)
    # Title priority: session_name → VoC heading → "Untitled Session"
    session_name = meta.get("session_name") or ""
    voc = state.get("voc_data", "")
    voc_heading = ""
    for line in voc.split("\n"):
        if line.startswith("## ") and "Perspective" not in line:
            voc_heading = line.replace("## ", "").strip()
            break
    title = session_name.strip() or voc_heading or "Untitled Session"

    _rect(s, 2.5, 1.3, 8.3, 0.5, HPE_GREEN_D)
    _tb(s, 2.5, 1.32, 8.3, 0.5, "BMI WORKFLOW REPORT", FH, 14, BG_DARK, True, PP_ALIGN.CENTER)
    _tb(s, 0.8, 2.5, 11.7, 0.8, title, FH, 34, WHITE, True, PP_ALIGN.CENTER)

    patterns = state.get("selected_patterns", [])
    direction = state.get("pattern_direction", "")
    subtitle = f"Direction: {direction.upper()}  \u00b7  {', '.join(patterns)}" if patterns else ""
    if subtitle:
        _tb(s, 0.8, 3.5, 11.7, 0.5, subtitle, FH, 18, HPE_GREEN, True, PP_ALIGN.CENTER)

    info_lines = [
        f"Session: {meta.get('session_id', '')}",
        f"Date: {meta.get('created', 'N/A')}  \u00b7  Backend: {meta.get('llm_backend', '')}",
        f"Status: {meta.get('status', '')}",
    ]
    _multi_tb(s, 2.5, 4.5, 8.3, 1.2, info_lines, FB, 12, MUTED, align=PP_ALIGN.CENTER, sp=6)

    _place_illustration(s, illus, SLIDE_ILLUSTRATIONS.get("cover", []),
                        left=11.5, top=1.2, size=0.9)
    _footer(s)
    _notes(s, (
        f"BMI WORKFLOW REPORT\n"
        f"Title: {title}\n"
        f"Direction: {direction}  |  Patterns: {', '.join(patterns)}\n"
        f"Session ID: {meta.get('session_id', '')}\n"
        f"Date: {meta.get('created', 'N/A')}\n"
        f"Backend: {meta.get('llm_backend', '')}\n"
        f"Status: {meta.get('status', '')}\n"
    ))
    return s


def _slide_signals(prs, signals, illus):
    """Slide: Detected Signals — full content cards."""
    s = _slide(prs)
    _section_label(s, 0.3, "STEP 1 \u2014 SOC RADAR")
    _section_title(s, 0.3, "Detected Signals")
    _place_illustration(s, illus, SLIDE_ILLUSTRATIONS.get("signals", []))

    max_cards = min(len(signals), 5)
    if max_cards == 0:
        _tb(s, 1, 3, 11, 1, "No signals detected.", FB, 12, MUTED)
        _footer(s)
        return s

    if max_cards <= 3:
        cols, rows_n = max_cards, 1
    else:
        cols, rows_n = 3, 2

    cw = 3.7
    ch = 2.4
    gx, gy = 0.35, 0.3
    total_w = cols * cw + (cols - 1) * gx
    sx = (13.33 - total_w) / 2
    sy = 1.5

    zone_colors = {
        "Business Model Anomaly": HPE_ORANGE,
        "Nonconsumption": HPE_GREEN,
        "Overserved Customers": HPE_BLUE,
        "Low-End Foothold": HPE_PURPLE,
        "New-Market Foothold": HPE_CYAN,
        "Enabling Technology": HPE_YELLOW,
        "Regulatory / Policy Shift": HPE_ORANGE,
    }

    for idx, sig in enumerate(signals[:max_cards]):
        col = idx % cols
        row = idx // cols
        x = sx + col * (cw + gx)
        y = sy + row * (ch + gy)
        zone = sig.get("zone", "")
        accent = zone_colors.get(zone, HPE_GREEN)

        _rect(s, x, y, cw, ch, CARD_BG, CARD_BORDER)
        _bar(s, x, y, cw, 0.05, accent)
        _tb(s, x + 0.15, y + 0.1, cw - 0.3, 0.25, zone, FH, 11, accent, True)
        # Full signal text
        _tb(s, x + 0.15, y + 0.4, cw - 0.3, 1.0,
            _clean(sig.get("signal", "")), FB, 11, WHITE)
        # First evidence item — full
        evidence = sig.get("evidence", [])
        if evidence:
            _tb(s, x + 0.15, y + 1.5, cw - 0.3, 0.8,
                _clean(evidence[0]), FB, 11, MUTED)

    _footer(s)

    notes_lines = ["DETECTED SIGNALS\n"]
    for i, sig in enumerate(signals):
        notes_lines.append(f"Signal {i+1}: {sig.get('signal', '')}")
        notes_lines.append(f"  Zone: {sig.get('zone', '')}")
        notes_lines.append(f"  Source: {sig.get('source', '')}")
        notes_lines.append(f"  Observable Behavior: {sig.get('observable_behavior', '')}")
        for j, ev in enumerate(sig.get("evidence", [])):
            notes_lines.append(f"  Evidence {j+1}: {ev}")
        notes_lines.append("")
    _notes(s, "\n".join(notes_lines))
    return s


def _slide_priority(prs, matrix, illus):
    """Slide: Priority Matrix — full text rows."""
    s = _slide(prs)
    _section_label(s, 0.3, "STEP 1 \u2014 SOC RADAR")
    _section_title(s, 0.3, "Priority Matrix")
    _place_illustration(s, illus, SLIDE_ILLUSTRATIONS.get("priority", []))

    if not matrix:
        _tb(s, 1, 3, 11, 1, "No priority matrix available.", FB, 12, MUTED)
        _footer(s)
        return s

    sorted_matrix = sorted(matrix, key=lambda x: x.get("score", 0), reverse=True)
    y = 1.5
    for idx, p in enumerate(sorted_matrix[:5]):
        tier = p.get("tier", "Monitor")
        accent = TIER_COLORS.get(tier, MID_GRAY)
        score = p.get("score", 0)

        _rect(s, 0.7, y, 11.9, 0.95, CARD_BG, CARD_BORDER)
        _bar(s, 0.7, y, 0.08, 0.95, accent)

        _rect(s, 0.95, y + 0.12, 0.55, 0.55, accent, radius=True)
        _tb(s, 0.95, y + 0.14, 0.55, 0.55, str(score), FH, 14, BG_DARK, True, PP_ALIGN.CENTER)

        # Full signal and rationale
        _tb(s, 1.7, y + 0.05, 8.3, 0.4,
            _clean(p.get("signal", "")), FB, 12, WHITE, True)
        _tb(s, 1.7, y + 0.45, 8.3, 0.45,
            _clean(p.get("rationale", "")), FB, 11, MUTED)

        _rect(s, 10.5, y + 0.2, 1.8, 0.4, accent)
        _tb(s, 10.5, y + 0.22, 1.8, 0.4, tier.upper(), FH, 11, BG_DARK, True, PP_ALIGN.CENTER)

        y += 1.05

    _footer(s)

    notes_lines = ["PRIORITY MATRIX\n"]
    for p in sorted_matrix:
        notes_lines.append(f"Signal: {p.get('signal', '')}")
        notes_lines.append(f"  Impact: {p.get('impact', '')}  Speed: {p.get('speed', '')}  Score: {p.get('score', '')}")
        notes_lines.append(f"  Tier: {p.get('tier', '')}")
        notes_lines.append(f"  Rationale: {p.get('rationale', '')}")
        notes_lines.append("")
    _notes(s, "\n".join(notes_lines))
    return s


def _slide_pattern(prs, state, illus):
    """Slide: Pattern Selection — full recommendation."""
    s = _slide(prs)
    _section_label(s, 0.3, "STEP 2 \u2014 PATTERN SELECTION")
    _section_title(s, 0.3, "Innovation Direction & Patterns")
    _place_illustration(s, illus, SLIDE_ILLUSTRATIONS.get("pattern", []))

    direction = state.get("pattern_direction", "").upper()
    patterns = state.get("selected_patterns", [])

    _rect(s, 4.2, 1.5, 5.0, 0.5, HPE_GREEN_D)
    _tb(s, 4.2, 1.52, 5.0, 0.5, f"DIRECTION: {direction}", FH, 14, BG_DARK, True, PP_ALIGN.CENTER)

    pcw = 3.8
    total_p = len(patterns) * pcw + (len(patterns) - 1) * 0.4 if patterns else 0
    px = (13.33 - total_p) / 2
    for j, pat in enumerate(patterns):
        x = px + j * (pcw + 0.4)
        _rect(s, x, 2.3, pcw, 0.5, CARD_BG, HPE_GREEN)
        _tb(s, x, 2.32, pcw, 0.5, pat, FH, 12, WHITE, True, PP_ALIGN.CENTER)

    # Full agent recommendation
    rec = state.get("agent_recommendation", "")
    if rec:
        rec_clean = _clean(rec.strip().replace("\n", " "))
        _multi_tb(s, 0.7, 3.2, 11.9, 3.5, [rec_clean], FB, 11, MUTED, sp=4)

    _footer(s)
    _notes(s, (
        f"PATTERN SELECTION\n"
        f"Direction: {direction}\n"
        f"Patterns: {', '.join(patterns)}\n\n"
        f"Full Recommendation:\n{rec}\n"
    ))
    return s


def _slide_customer_profile(prs, profile_md, illus):
    """Slide: Customer Profile — table per Jobs/Pains/Gains."""
    s = _slide(prs)
    _section_label(s, 0.3, "STEP 3 \u2014 EMPATHIZE")
    _section_title(s, 0.3, "Customer Empathy Profile")
    _place_illustration(s, illus, SLIDE_ILLUSTRATIONS.get("profile", []))

    seg_match = re.search(r"### Customer Segment\n(.+?)(?:\n###|\Z)", profile_md, re.DOTALL)
    segment = seg_match.group(1).strip().split("\n")[0] if seg_match else ""
    if segment:
        _tb(s, 0.7, 1.2, 11.0, 0.35, _clean(segment), FB, 11, MUTED)

    sections = {
        "Jobs": {"accent": HPE_GREEN, "label": "CUSTOMER JOBS",
                 "val_col": "Job", "imp_col": "Importance"},
        "Pains": {"accent": HPE_ORANGE, "label": "CUSTOMER PAINS",
                  "val_col": "Pain", "imp_col": "Severity"},
        "Gains": {"accent": HPE_BLUE, "label": "CUSTOMER GAINS",
                  "val_col": "Gain", "imp_col": "Relevance"},
    }

    notes_lines = ["CUSTOMER EMPATHY PROFILE\n"]
    if segment:
        notes_lines.append(f"Segment: {segment}\n")

    y = 1.65
    for key, cfg in sections.items():
        table_match = re.search(
            rf"### Customer {key}\n(\|.+?)(?=\n###|\n---|\Z)", profile_md, re.DOTALL
        )
        if not table_match:
            continue
        rows = _parse_md_table(table_match.group(1))
        if not rows:
            continue

        _tb(s, 0.7, y, 4, 0.25, cfg["label"], FH, 11, cfg["accent"], True)
        y += 0.28

        notes_lines.append(f"\n{cfg['label']}:")
        tbl_rows = []
        for row in rows:
            type_val = row.get("Type", "")
            content = row.get(cfg["val_col"], "")
            importance = row.get(cfg["imp_col"], "")
            notes_lines.append(f"  [{type_val}] {content} \u2014 {importance}")
            tbl_rows.append([type_val, _clean(content), importance])

        _styled_table(s, 0.7, y, 11.5,
                      [1.2, 8.0, 2.3],
                      ["Type", cfg["val_col"], cfg["imp_col"]],
                      tbl_rows[:6],
                      accent=cfg["accent"], font_sz=11)
        y += 0.35 * (min(len(tbl_rows), 6) + 1) + 0.15

    _footer(s)
    _notes(s, "\n".join(notes_lines))
    return s


def _slide_value_driver_tree(prs, vdt_md, illus):
    """Slide: Value Driver Tree — full table."""
    s = _slide(prs)
    _section_label(s, 0.3, "STEP 4 \u2014 MEASURE")
    _section_title(s, 0.3, "Value Driver Tree")
    _place_illustration(s, illus, SLIDE_ILLUSTRATIONS.get("vdt", []))

    outcome_match = re.search(r"### Customer Business Outcome\n(.+?)(?:\n###|\Z)", vdt_md, re.DOTALL)
    outcome = outcome_match.group(1).strip().split("\n")[0] if outcome_match else ""

    notes_lines = ["VALUE DRIVER TREE\n"]
    if outcome:
        _rect(s, 0.7, 1.35, 11.5, 0.4, CARD_BG2)
        _tb(s, 0.9, 1.37, 11.3, 0.4, _clean(outcome), FB, 12, WHITE)
        notes_lines.append(f"Business Outcome: {outcome}\n")

    table_match = re.search(
        r"### Key Deliverables.*?\n(\|.+?)(?=\n###|\n---|\Z)", vdt_md, re.DOTALL
    )
    if table_match:
        rows = _parse_md_table(table_match.group(1))
        notes_lines.append("Key Deliverables:")
        tbl_rows = []
        for row in rows:
            deliv = _clean(row.get("Key Deliverable", ""))
            measure = _clean(row.get("Success Measure", ""))
            baseline = row.get("Baseline", "")
            target = row.get("Target", "")
            driver = row.get("Driver Type", "")
            notes_lines.append(f"\n  Deliverable: {deliv}")
            notes_lines.append(f"  Measure: {measure}")
            notes_lines.append(f"  Baseline: {baseline}")
            notes_lines.append(f"  Target: {target}")
            notes_lines.append(f"  Driver: {driver}")
            tbl_rows.append([deliv, measure, baseline, target, driver])

        _styled_table(s, 0.7, 1.95, 11.5,
                      [3.0, 3.5, 1.8, 1.8, 1.4],
                      ["Deliverable", "Success Measure", "Baseline", "Target", "Driver"],
                      tbl_rows,
                      accent=HPE_GREEN, font_sz=11)

    _footer(s)
    _notes(s, "\n".join(notes_lines))
    return s


def _slide_actionable_insights(prs, insights_md, illus):
    """Slide 1: Actionable Insights — XYZ + Problem Statements."""
    s = _slide(prs)
    _section_label(s, 0.3, "STEP 4 \u2014 DEFINE")
    _section_title(s, 0.3, "Actionable Insights")
    _place_illustration(s, illus, SLIDE_ILLUSTRATIONS.get("insights", []))

    notes_lines = ["ACTIONABLE INSIGHTS\n"]

    xyz_match = re.search(r"\*\*(.+?)\*\* DOES \*\*(.+?)\*\* BECAUSE \*\*(.+?)\*\* BUT \*\*(.+?)\*\*",
                          insights_md)
    y = 1.4
    if xyz_match:
        who, does, because, but = xyz_match.groups()
        notes_lines.append(f"WHO: {who}")
        notes_lines.append(f"DOES: {does}")
        notes_lines.append(f"BECAUSE: {because}")
        notes_lines.append(f"BUT: {but}\n")

        _rect(s, 0.7, y, 11.5, 1.5, CARD_BG2)
        _tb(s, 0.9, y + 0.05, 11.1, 0.3, _clean(who), FH, 12, HPE_GREEN, True)
        _tb(s, 0.9, y + 0.35, 11.1, 0.3, f"DOES: {_clean(does)}", FB, 11, WHITE)
        _tb(s, 0.9, y + 0.65, 11.1, 0.3, f"BECAUSE: {_clean(because)}", FB, 11, MUTED)
        _tb(s, 0.9, y + 0.95, 11.1, 0.5, f"BUT: {_clean(but)}", FB, 11, HPE_ORANGE)
        y += 1.65

    prob_match = re.search(r"### Problem Statements\n(\|.+?)(?=\n---|\Z)", insights_md, re.DOTALL)
    if prob_match:
        rows = _parse_md_table(prob_match.group(1))
        notes_lines.append("PROBLEM STATEMENTS:")
        _tb(s, 0.7, y, 4, 0.25, "PROBLEM STATEMENTS", FH, 11, HPE_ORANGE, True)
        y += 0.3

        tbl_rows = []
        for row in rows:
            num = row.get("#", "")
            stmt = _clean(row.get("Problem Statement", ""))
            priority = row.get("Priority", "")
            jobs = _clean(row.get("Jobs Affected", ""))
            pains = _clean(row.get("Pains Addressed", ""))
            notes_lines.append(f"\n  #{num} [{priority}] {stmt}")
            notes_lines.append(f"  Jobs: {jobs}")
            notes_lines.append(f"  Pains: {pains}")
            tbl_rows.append([num, stmt, priority, jobs, pains])

        _styled_table(s, 0.7, y, 11.5,
                      [0.4, 4.0, 0.9, 3.0, 3.0],
                      ["#", "Problem Statement", "Priority", "Jobs Affected", "Pains Addressed"],
                      tbl_rows,
                      accent=HPE_ORANGE, font_sz=11)

    _footer(s)
    _notes(s, "\n".join(notes_lines))
    return s


def _slide_actionable_insights_2(prs, insights_md, illus):
    """Slide 2: Value Chain Assessment + Customer Journey Friction Points."""
    vc_match = re.search(r"### Value Chain Assessment\n(\|.+?)(?=\n###|\Z)", insights_md, re.DOTALL)
    fr_match = re.search(r"### Customer Journey Friction Points\n(\|.+?)(?=\n###|\Z)", insights_md, re.DOTALL)
    if not vc_match and not fr_match:
        return None

    s = _slide(prs)
    _section_label(s, 0.3, "STEP 4 \u2014 DEFINE")
    _section_title(s, 0.3, "Value Chain & Journey Friction")
    _place_illustration(s, illus, SLIDE_ILLUSTRATIONS.get("insights2", []))

    notes_lines = ["VALUE CHAIN & JOURNEY FRICTION\n"]
    y = 1.4

    if vc_match:
        rows = _parse_md_table(vc_match.group(1))
        _tb(s, 0.7, y, 6, 0.25, "VALUE CHAIN ASSESSMENT", FH, 11, HPE_BLUE, True)
        y += 0.3
        notes_lines.append("VALUE CHAIN ASSESSMENT:")
        tbl_rows = []
        for row in rows:
            activity = _clean(row.get("Activity", ""))
            role = _clean(row.get("Role in Value Creation", ""))
            weak = row.get("Weak Link?", "")
            impact = _clean(row.get("Impact on Customer", ""))
            notes_lines.append(f"  {activity}: {role} (Weak: {weak}) Impact: {impact}")
            tbl_rows.append([activity, role, weak, impact])

        _styled_table(s, 0.7, y, 11.5,
                      [2.5, 3.5, 1.0, 4.5],
                      ["Activity", "Role in Value Creation", "Weak?", "Impact on Customer"],
                      tbl_rows[:5],
                      accent=HPE_BLUE, font_sz=11)
        y += 0.35 * (min(len(tbl_rows), 5) + 1) + 0.3

    if fr_match:
        rows = _parse_md_table(fr_match.group(1))
        _tb(s, 0.7, y, 6, 0.25, "CUSTOMER JOURNEY FRICTION", FH, 11, HPE_PURPLE, True)
        y += 0.3
        notes_lines.append("\nCUSTOMER JOURNEY FRICTION:")
        tbl_rows = []
        for row in rows:
            phase = row.get("Journey Phase", "")
            touchpoint = _clean(row.get("Touchpoint", ""))
            exp = _clean(row.get("Customer Experience", ""))
            friction = row.get("Friction Type", "")
            opp = _clean(row.get("Opportunity", ""))
            notes_lines.append(f"  {phase}: {touchpoint} [{friction}] -> {opp}")
            tbl_rows.append([phase, touchpoint, friction, opp])

        _styled_table(s, 0.7, y, 11.5,
                      [1.8, 3.0, 1.5, 5.2],
                      ["Phase", "Touchpoint", "Friction", "Opportunity"],
                      tbl_rows[:5],
                      accent=HPE_PURPLE, font_sz=11)

    _footer(s)
    _notes(s, "\n".join(notes_lines))
    return s


def _slide_vpc(prs, vpc_md, illus):
    """Slide: Value Proposition Canvas — table per section."""
    s = _slide(prs)
    _section_label(s, 0.3, "STEP 5 \u2014 VALUE PROPOSITION")
    _section_title(s, 0.3, "Value Proposition Canvas")
    _place_illustration(s, illus, SLIDE_ILLUSTRATIONS.get("vpc", []))

    card_defs = [
        ("Products & Services", "PRODUCTS & SERVICES", HPE_GREEN,
         r"#### Products & Services\n(\|.+?)(?=\n####|\n###|\Z)",
         "Product/Service", "Relevance"),
        ("Pain Relievers", "PAIN RELIEVERS", HPE_ORANGE,
         r"#### Pain Relievers\n(\|.+?)(?=\n####|\n###|\Z)",
         "Pain Reliever", "Pain Addressed"),
        ("Gain Creators", "GAIN CREATORS", HPE_BLUE,
         r"#### Gain Creators\n(\|.+?)(?=\n####|\n###|\Z)",
         "Gain Creator", "Gain Addressed"),
    ]

    notes_lines = ["VALUE PROPOSITION CANVAS\n"]
    y = 1.35

    for key, label, accent, pattern, val_col, detail_col in card_defs:
        match = re.search(pattern, vpc_md, re.DOTALL)
        if not match:
            continue
        rows = _parse_md_table(match.group(1))
        if not rows:
            continue

        _tb(s, 0.7, y, 6, 0.25, label, FH, 11, accent, True)
        y += 0.28

        notes_lines.append(f"\n{key}:")
        tbl_rows = []
        for row in rows:
            content = _clean(row.get(val_col, ""))
            relevance = row.get("Relevance", "")
            detail = _clean(row.get(detail_col, ""))
            notes_lines.append(f"  [{relevance}] {content}")
            if detail and detail_col != "Relevance":
                notes_lines.append(f"    {detail_col}: {detail}")
            tbl_rows.append([content, relevance, detail])

        headers = [val_col, "Relevance", detail_col] if detail_col != "Relevance" else [val_col, "Relevance"]
        widths = [5.0, 1.5, 5.0] if detail_col != "Relevance" else [8.0, 3.5]
        tbl_data = tbl_rows if detail_col != "Relevance" else [[r[0], r[1]] for r in tbl_rows]

        _styled_table(s, 0.7, y, 11.5, widths, headers, tbl_data[:4],
                      accent=accent, font_sz=11)
        y += 0.35 * (min(len(tbl_data), 4) + 1) + 0.15

    # Ad-lib prototype
    adlib_match = re.search(r"> \*\*OUR\*\*(.+?)(?=Context anchor|\Z)", vpc_md, re.DOTALL)
    if adlib_match:
        adlib_full = _clean(adlib_match.group(1).strip().replace("> ", "").replace("\n", " "))
        notes_lines.append(f"\n\nAD-LIB PROTOTYPES:\n{adlib_full}")
        if y < 6.2:
            _rect(s, 0.7, y + 0.05, 11.5, 0.55, CARD_BG2)
            _tb(s, 0.85, y + 0.07, 1.0, 0.25, "AD-LIB", FH, 11, HPE_GREEN, True)
            _tb(s, 0.85, y + 0.3, 11.2, 0.3, adlib_full, FB, 11, MUTED)

    _footer(s)
    _notes(s, "\n".join(notes_lines))
    return s


def _slide_bmc(prs, bmc_md, illus):
    """Slide: Business Model Canvas — unified table."""
    s = _slide(prs)
    _section_label(s, 0.3, "STEP 6 \u2014 DESIGN")
    _section_title(s, 0.3, "Business Model Canvas")
    _place_illustration(s, illus, SLIDE_ILLUSTRATIONS.get("bmc", []))

    notes_lines = ["BUSINESS MODEL CANVAS\n"]

    bmc_dims = [
        ("Desirability", HPE_GREEN),
        ("Feasibility", HPE_BLUE),
        ("Viability", HPE_PURPLE),
    ]

    all_rows = []
    for dim, accent in bmc_dims:
        dim_match = re.search(rf"### {dim}\n(\|.+?)(?=\n###|\n---|\Z)", bmc_md, re.DOTALL)
        if dim_match:
            rows = _parse_md_table(dim_match.group(1))
            notes_lines.append(f"\n{dim.upper()}:")
            for row in rows:
                block = _clean(row.get("Building Block", ""))
                desc = _clean(row.get("Description", ""))
                notes_lines.append(f"  {block}: {desc}")
                all_rows.append([dim, block, desc])

    if all_rows:
        _styled_table(s, 0.7, 1.4, 11.5,
                      [1.8, 2.5, 7.2],
                      ["Dimension", "Building Block", "Description"],
                      all_rows,
                      accent=HPE_GREEN, font_sz=11)

    _footer(s)
    _notes(s, "\n".join(notes_lines))
    return s


def _slide_fit(prs, fit_md, illus):
    """Slide: Fit Assessment — table per fit type."""
    s = _slide(prs)
    _section_label(s, 0.3, "STEP 6 \u2014 FIT ASSESSMENT")
    _section_title(s, 0.3, "Problem-Solution / Product-Market / BM Fit")
    _place_illustration(s, illus, SLIDE_ILLUSTRATIONS.get("fit", []))

    fit_types = [
        ("Problem-Solution Fit", HPE_GREEN, r"### Problem-Solution Fit\n(\|.+?)(?=\n###|\Z)"),
        ("Product-Market Fit", HPE_BLUE, r"### Product-Market Fit Status\n(\|.+?)(?=\n###|\Z)"),
        ("Business Model Fit", HPE_PURPLE, r"### Business Model Fit Status\n(\|.+?)(?=\n###|\Z)"),
    ]

    notes_lines = ["FIT ASSESSMENT\n"]
    y = 1.35

    for title, accent, pattern in fit_types:
        match = re.search(pattern, fit_md, re.DOTALL)
        if not match:
            continue
        rows = _parse_md_table(match.group(1))
        if not rows:
            continue

        _tb(s, 0.7, y, 6, 0.25, title.upper(), FH, 11, accent, True)
        y += 0.28

        notes_lines.append(f"\n{title.upper()}:")
        tbl_rows = []
        for row in rows:
            need = _clean(row.get("Customer Need (Job/Pain/Gain)") or row.get("Criterion")
                          or row.get("Dimension") or "")
            status = row.get("Fit?") or row.get("Status") or ""
            evidence = _clean(row.get("Evidence") or row.get("Mapped Value Proposition Element") or "")
            importance = row.get("Importance to Customer", "")
            notes_lines.append(f"  {need}: {status} \u2014 {evidence}")
            if importance:
                notes_lines.append(f"    Importance: {importance}")
            tbl_rows.append([need, status, evidence])

        _styled_table(s, 0.7, y, 11.5,
                      [3.5, 1.2, 6.8],
                      ["Need / Criterion", "Status", "Evidence"],
                      tbl_rows[:5],
                      accent=accent, font_sz=11)
        y += 0.35 * (min(len(tbl_rows), 5) + 1) + 0.15

    _footer(s)
    _notes(s, "\n".join(notes_lines))
    return s


def _slide_assumptions(prs, assumptions_md, illus):
    """Slide: Assumptions Map — TABLE layout."""
    s = _slide(prs)
    _section_label(s, 0.3, "STEP 7 \u2014 DE-RISK")
    _section_title(s, 0.3, "Assumption Map")
    _place_illustration(s, illus, SLIDE_ILLUSTRATIONS.get("assumptions", []))

    notes_lines = ["ASSUMPTION MAP\n"]

    quad_match = re.search(r"## Importance.*?Evidence Map\n(\|.+?)(?=\n##|\Z)", assumptions_md, re.DOTALL)
    if quad_match:
        rows = _parse_md_table(quad_match.group(1))
        notes_lines.append("Importance vs. Evidence Map:")
        tbl_rows = []
        for row in rows:
            assumption = _clean(row.get("Assumption", ""))
            category = row.get("Category", "")
            quadrant = row.get("Suggested Quadrant", "") or row.get("Quadrant", "")
            importance = row.get("Importance", "")
            evidence = row.get("Evidence", "")
            notes_lines.append(f"  [{category}] {assumption} -> {quadrant}")
            tbl_rows.append([assumption, category, importance, evidence, quadrant])

        _styled_table(s, 0.7, 1.4, 11.5,
                      [4.0, 1.5, 1.2, 1.2, 2.6],
                      ["Assumption", "Category", "Imp.", "Evid.", "Quadrant"],
                      tbl_rows,
                      accent=HPE_ORANGE, font_sz=11)
    else:
        # Fallback: dimension-level tables
        all_rows = []
        for dim in ["Desirability", "Viability", "Feasibility"]:
            dim_match = re.search(rf"## {dim}\n(\|.+?)(?=\n##|\Z)", assumptions_md, re.DOTALL)
            if dim_match:
                rows = _parse_md_table(dim_match.group(1))
                notes_lines.append(f"\n{dim.upper()}:")
                for row in rows:
                    assumption = _clean(row.get("Assumption", ""))
                    category = row.get("Category", "")
                    notes_lines.append(f"  [{category}] {assumption}")
                    all_rows.append([dim, category, assumption])

        if all_rows:
            _styled_table(s, 0.7, 1.4, 11.5,
                          [1.8, 1.5, 8.2],
                          ["Dimension", "Category", "Assumption"],
                          all_rows,
                          accent=HPE_ORANGE, font_sz=11)

    _footer(s)
    _notes(s, "\n".join(notes_lines))
    return s


def _slide_dvf_tensions(prs, assumptions_md, illus):
    """Slide: DVF Tensions — cards with full text."""
    tension_match = re.search(r"## DVF Tensions\n(\|.+?)(?=\n##|\Z)", assumptions_md, re.DOTALL)
    if not tension_match:
        return None
    rows = _parse_md_table(tension_match.group(1))
    if not rows:
        return None

    s = _slide(prs)
    _section_label(s, 0.3, "STEP 7 \u2014 DE-RISK")
    _section_title(s, 0.3, "DVF Tensions")
    _place_illustration(s, illus, SLIDE_ILLUSTRATIONS.get("dvf", []))

    notes_lines = ["DVF TENSIONS\n"]
    y = 1.7
    for idx, row in enumerate(rows[:4]):
        tension = _clean(row.get("Tension", ""))
        categories = row.get("Categories", "")
        notes_lines.append(f"\nTension {idx+1}: {tension}")
        notes_lines.append(f"  Categories: {categories}")

        _rect(s, 0.7, y, 11.9, 1.1, CARD_BG, CARD_BORDER)
        _bar(s, 0.7, y, 0.08, 1.1, HPE_ORANGE)
        _tb(s, 1.0, y + 0.05, 11.0, 0.6, tension, FB, 11, WHITE)
        _tb(s, 1.0, y + 0.65, 11.0, 0.35, categories, FH, 11, HPE_ORANGE, True)
        y += 1.3

    _footer(s)
    _notes(s, "\n".join(notes_lines))
    return s


def _slide_experiments(prs, experiment_cards, illus):
    """Slide: Experiment Selection — full content cards."""
    s = _slide(prs)
    _section_label(s, 0.3, "STEP 8 \u2014 TEST")
    _section_title(s, 0.3, "Experiment Selection")
    _place_illustration(s, illus, SLIDE_ILLUSTRATIONS.get("experiments", []))

    if not experiment_cards:
        _tb(s, 1, 3, 11, 1, "No experiments selected.", FB, 12, MUTED)
        _footer(s)
        return s

    by_cat: dict[str, list[dict]] = {}
    for card in experiment_cards:
        cat = card.get("category", "Unknown")
        by_cat.setdefault(cat, []).append(card)

    cols = list(by_cat.keys())
    n_cols = len(cols)
    cw = min(3.8, (12.0 - (n_cols - 1) * 0.35) / n_cols)
    gx = 0.35
    total_w = n_cols * cw + (n_cols - 1) * gx
    sx = (13.33 - total_w) / 2
    sy = 1.6
    ch = 4.8

    notes_lines = ["EXPERIMENT SELECTION\n"]

    for j, cat in enumerate(cols):
        x = sx + j * (cw + gx)
        accent = CAT_COLORS.get(cat, HPE_GREEN)
        cards_in_cat = by_cat[cat]

        _rect(s, x, sy, cw, ch, CARD_BG, CARD_BORDER)
        _bar(s, x, sy, cw, 0.05, accent)
        _tb(s, x + 0.15, sy + 0.12, cw - 0.3, 0.35, cat.upper(), FH, 12, accent, True)

        notes_lines.append(f"\n{cat.upper()}:")
        lines = []
        for c in cards_in_cat:
            assumption = _clean(c.get("assumption", ""))
            card_name = c.get("card_name", "")
            strength = c.get("evidence_strength", "")
            path = c.get("evidence_path", [])
            path_str = " \u2192 ".join(p.get("card_name", "") for p in path) if path else ""

            notes_lines.append(f"  Card: {card_name} [{strength}]")
            notes_lines.append(f"    Assumption: {assumption}")
            if path_str:
                notes_lines.append(f"    Path: {path_str}")

            lines.append(f"\u25B6 {card_name}")
            lines.append(f"  {assumption}")

        if lines:
            _multi_tb(s, x + 0.15, sy + 0.55, cw - 0.3, ch - 0.7,
                      lines, FB, 11, MUTED, sp=3)

    _footer(s)
    _notes(s, "\n".join(notes_lines))
    return s


def _slide_experiment_worksheets(prs, worksheets_md, illus):
    """Slide(s): Experiment Worksheets — full content cards."""
    chunks = re.split(r"## Experiment Worksheet\b", worksheets_md)
    chunks = [c.strip() for c in chunks if c.strip()]

    if not chunks:
        return []

    slides = []
    for idx, chunk in enumerate(chunks[:4]):
        s = _slide(prs)
        _section_label(s, 0.3, "STEP 8 \u2014 EXPERIMENT WORKSHEET")
        _place_illustration(s, illus, SLIDE_ILLUSTRATIONS.get("worksheets", []))

        card_match = re.search(r"\*\*Experiment card:\*\* (.+)", chunk)
        cat_match = re.search(r"\*\*Category:\*\* (.+)", chunk)
        assumption_match = re.search(r"\*\*Assumption statement:\*\* (.+)", chunk)
        primary_metric_match = re.search(r"\*\*Primary metric:\*\* (.+)", chunk)
        success_match = re.search(r"\*\*Success looks like:\*\* (.+)", chunk)
        failure_match = re.search(r"\*\*Failure looks like:\*\* (.+)", chunk)
        timebox_match = re.search(r"\*\*Timebox:\*\* (.+)", chunk)
        sample_match = re.search(r"\*\*Sample size target:\*\* (.+)", chunk)

        card_name = card_match.group(1) if card_match else f"Experiment {idx + 1}"
        category = cat_match.group(1) if cat_match else ""
        assumption = _clean(assumption_match.group(1)) if assumption_match else ""
        accent = CAT_COLORS.get(category, HPE_GREEN)

        _section_title(s, 0.3, f"{card_name}")

        _rect(s, 0.7, 1.5, 11.9, 5.0, CARD_BG, CARD_BORDER)
        _bar(s, 0.7, 1.5, 11.9, 0.05, accent)

        _tb(s, 0.9, 1.65, 1.8, 0.3, "ASSUMPTION", FH, 11, accent, True)
        _tb(s, 0.9, 1.95, 11.3, 0.5, assumption, FB, 11, WHITE)

        col1_x, col2_x = 0.9, 6.8
        col_w = 5.5

        _tb(s, col1_x, 2.5, col_w, 0.3, "EXPERIMENT DESIGN", FH, 11, accent, True)
        design_lines = []
        if primary_metric_match:
            design_lines.append(f"Metric: {_clean(primary_metric_match.group(1))}")
        if sample_match:
            design_lines.append(f"Sample: {sample_match.group(1)}")
        if timebox_match:
            design_lines.append(f"Timebox: {timebox_match.group(1)}")
        if design_lines:
            _multi_tb(s, col1_x, 2.85, col_w, 1.5, design_lines, FB, 11, MUTED, sp=3, bullet=True)

        _tb(s, col2_x, 2.5, col_w, 0.3, "DECISION CRITERIA", FH, 11, accent, True)
        criteria_lines = []
        if success_match:
            criteria_lines.append(f"\u2705 {_clean(success_match.group(1))}")
        if failure_match:
            criteria_lines.append(f"\u274C {_clean(failure_match.group(1))}")
        if criteria_lines:
            _multi_tb(s, col2_x, 2.85, col_w, 1.5, criteria_lines, FB, 11, MUTED, sp=3)

        exec_match = re.search(r"### Execution Plan\n((?:\d+\..*\n?)+)", chunk)
        steps = []
        if exec_match:
            _tb(s, col1_x, 4.3, 11, 0.3, "EXECUTION PLAN", FH, 11, accent, True)
            steps = [l.strip() for l in exec_match.group(1).strip().split("\n") if l.strip()]
            step_full = [_clean(st) for st in steps[:4]]
            _multi_tb(s, col1_x, 4.65, 11, 1.2, step_full, FB, 11, MUTED, sp=3)

        seq_pos_match = re.search(r"\*\*If signal is positive, move next to:\*\* (.+)", chunk)
        seq_mix_match = re.search(r"\*\*If signal is weak or mixed, move next to:\*\* (.+)", chunk)
        if seq_pos_match:
            _tb(s, col1_x, 5.8, 11, 0.3, "SEQUENCING", FH, 11, accent, True)
            seq_lines = [f"\u2192 Positive: {seq_pos_match.group(1)}"]
            if seq_mix_match:
                seq_lines.append(f"\u2192 Mixed: {seq_mix_match.group(1)}")
            _multi_tb(s, col1_x, 6.1, 11, 0.5, seq_lines, FB, 11, MUTED, sp=3)

        _footer(s)

        notes_lines = [
            f"EXPERIMENT WORKSHEET: {card_name}\n",
            f"Category: {category}",
            f"Assumption: {assumption}",
        ]
        if primary_metric_match:
            notes_lines.append(f"Primary Metric: {primary_metric_match.group(1)}")
        if sample_match:
            notes_lines.append(f"Sample Size: {sample_match.group(1)}")
        if timebox_match:
            notes_lines.append(f"Timebox: {timebox_match.group(1)}")
        if success_match:
            notes_lines.append(f"\nSuccess: {success_match.group(1)}")
        if failure_match:
            notes_lines.append(f"Failure: {failure_match.group(1)}")
        if exec_match:
            notes_lines.append(f"\nExecution Plan:")
            for step in steps:
                notes_lines.append(f"  {step}")
        if seq_pos_match:
            notes_lines.append(f"\nSequencing:")
            notes_lines.append(f"  If positive: {seq_pos_match.group(1)}")
        if seq_mix_match:
            notes_lines.append(f"  If mixed: {seq_mix_match.group(1)}")

        _notes(s, "\n".join(notes_lines))
        slides.append(s)

    return slides


# ── Main ────────────────────────────────────────────────────────────────

def generate_report_pptx(
    state: dict[str, Any],
    meta: dict[str, Any],
    template_path: str = TEMPLATE,
    output_path: str | None = None,
    illustrations_path: str = ILLUSTRATIONS_PPTX,
) -> str:
    """Generate a PowerPoint report from workflow state."""
    prs = _init_prs(template_path)
    illus = _load_illustrations(illustrations_path)

    _slide_cover(prs, state, meta, illus)
    _slide_signals(prs, state.get("signals", []), illus)
    _slide_priority(prs, state.get("priority_matrix", []), illus)
    _slide_pattern(prs, state, illus)

    if state.get("customer_profile"):
        _slide_customer_profile(prs, state["customer_profile"], illus)
    if state.get("value_driver_tree"):
        _slide_value_driver_tree(prs, state["value_driver_tree"], illus)
    if state.get("actionable_insights"):
        _slide_actionable_insights(prs, state["actionable_insights"], illus)
        _slide_actionable_insights_2(prs, state["actionable_insights"], illus)
    if state.get("value_proposition_canvas"):
        _slide_vpc(prs, state["value_proposition_canvas"], illus)
    if state.get("business_model_canvas"):
        _slide_bmc(prs, state["business_model_canvas"], illus)
    if state.get("fit_assessment"):
        _slide_fit(prs, state["fit_assessment"], illus)
    if state.get("assumptions"):
        _slide_assumptions(prs, state["assumptions"], illus)
        _slide_dvf_tensions(prs, state["assumptions"], illus)

    _slide_experiments(prs, state.get("experiment_cards", []), illus)

    if state.get("experiment_worksheets"):
        _slide_experiment_worksheets(prs, state["experiment_worksheets"], illus)

    if output_path is None:
        output_path = f"reports/{meta.get('session_id', 'report')}_report.pptx"
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


# ── CLI ─────────────────────────────────────────────────────────────────

def main():
    import typer

    from backend.app.db.models import WorkflowRun
    from backend.app.db.session import SessionLocal

    app = typer.Typer()

    @app.command()
    def generate(
        session_id: str | None = typer.Option(None, "--session-id"),
        output: str | None = typer.Option(None, "--output", "-o"),
        template: str = typer.Option(TEMPLATE, "--template"),
    ):
        """Generate a PowerPoint report from a workflow session."""
        sess = SessionLocal()
        try:
            if session_id:
                run = sess.query(WorkflowRun).filter_by(session_id=session_id).first()
                if not run:
                    typer.echo(f"Error: session '{session_id}' not found.", err=True)
                    raise typer.Exit(code=1)
            else:
                run = (
                    sess.query(WorkflowRun)
                    .filter_by(status="completed")
                    .order_by(WorkflowRun.created_at.desc())
                    .first()
                )
                if not run:
                    run = sess.query(WorkflowRun).order_by(WorkflowRun.created_at.desc()).first()
                if not run:
                    typer.echo("Error: no workflow runs found.", err=True)
                    raise typer.Exit(code=1)
                typer.echo(f"Using most recent session: {run.session_id}")

            state = run.state_json if isinstance(run.state_json, dict) else {}
            meta = {
                "session_id": run.session_id,
                "session_name": run.session_name,
                "created": run.created_at.strftime("%Y-%m-%d %H:%M UTC") if run.created_at else "N/A",
                "status": run.status,
                "llm_backend": run.llm_backend,
                "input_type": run.input_type,
            }

            out = generate_report_pptx(state, meta, template_path=template, output_path=output)
            typer.echo(f"Report saved to {out}")
        finally:
            sess.close()

    app()


if __name__ == "__main__":
    main()
