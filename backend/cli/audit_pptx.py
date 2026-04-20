"""Audit a generated PPTX for quality: notes, fonts, truncation, illustrations."""
import sys
from pptx import Presentation


def audit(path: str) -> None:
    prs = Presentation(path)
    print(f"Total slides: {len(prs.slides)}\n")

    issues = []
    for i, slide in enumerate(prs.slides, 1):
        has_notes = False
        notes_len = 0
        try:
            nf = slide.notes_slide.notes_text_frame
            txt = nf.text.strip()
            if txt:
                has_notes = True
                notes_len = len(txt)
        except Exception:
            pass

        img_count = sum(1 for s in slide.shapes if s.shape_type == 13)

        min_font = None
        small_fonts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            pts = run.font.size / 12700
                            if min_font is None or pts < min_font:
                                min_font = pts
                            if pts < 16 and pts > 0:
                                small_fonts.append((shape.name, run.text[:40], pts))

        truncated = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t.endswith("...") and len(t) > 5:
                        truncated.append(t[:60])

        status = "Y" if has_notes else "N"
        mf = f"{min_font:.0f}" if min_font else "N/A"
        print(f"Slide {i:2d}: notes={status} ({notes_len:5d} chars) | imgs={img_count} | min_font={mf:>4}pt | trunc={len(truncated)}")

        if small_fonts:
            for name, text, sz in small_fonts[:3]:
                issues.append(f"  Slide {i}: {sz:.0f}pt in \"{text}\"")
        if truncated:
            for t in truncated[:2]:
                issues.append(f"  Slide {i} TRUNC: \"{t}\"")

    if issues:
        print(f"\n=== {len(issues)} Issues ===")
        for iss in issues:
            print(iss)
    else:
        print("\n=== No issues found! ===")


if __name__ == "__main__":
    audit(sys.argv[1] if len(sys.argv) > 1 else "reports/5fc8e37e-a11a-48ce-a3ac-c58199feca50_report.pptx")
