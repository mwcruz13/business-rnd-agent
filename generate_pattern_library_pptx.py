"""
Generate Strategyzer Pattern Library PowerPoint from the HPE Dark template.

Slides:
  1. Cover — 21 patterns overview
  2. INVENT overview — 9 patterns × 27 flavors summary grid
  3–11. One slide per INVENT pattern (with flavors)
  12. SHIFT overview — 12 patterns summary grid
  13–24. One slide per SHIFT pattern (strategic + reverse reflection)

Usage:
  python generate_pattern_library_pptx.py
"""
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# --- HPE Dark Theme Colors ---
BG_DARK     = RGBColor(0x29, 0x2D, 0x3A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
MUTED_TEXT  = RGBColor(0xB1, 0xB9, 0xBE)
MID_GRAY    = RGBColor(0x53, 0x5C, 0x66)
HPE_GREEN   = RGBColor(0x00, 0xE0, 0xAF)
HPE_PURPLE  = RGBColor(0x9B, 0x84, 0xFC)
HPE_BLUE    = RGBColor(0x65, 0xAE, 0xF9)
HPE_CYAN    = RGBColor(0x62, 0xE5, 0xF6)
HPE_ORANGE  = RGBColor(0xF2, 0x6B, 0x43)
CARD_BG     = RGBColor(0x33, 0x3A, 0x4A)

FONT_HEADING = "HPE Graphik Semibold"
FONT_BODY    = "HPE Graphik"

# --- Paths ---
SCRIPT_DIR = Path(__file__).resolve().parent
TEMPLATE   = str(SCRIPT_DIR / "hpe_dark_template.pptx")
DATA_FILE  = str(SCRIPT_DIR / "backend" / "app" / "patterns" / "strategyzer-pattern-library.json")
OUTPUT     = str(SCRIPT_DIR / "reports" / "Strategyzer_Pattern_Library.pptx")

# Category → accent color
INVENT_CAT_COLORS = {
    "Frontstage Disruption": HPE_GREEN,
    "Backstage Disruption": HPE_BLUE,
    "Profit Formula Disruption": HPE_PURPLE,
}

SHIFT_CAT_COLORS = {
    "Value Proposition Shifts": HPE_GREEN,
    "Frontstage Driven Shifts": HPE_CYAN,
    "Backstage Driven Shifts": HPE_BLUE,
    "Profit Formula Driven Shifts": HPE_PURPLE,
}


# ---- LOAD DATA ----
with open(DATA_FILE, encoding="utf-8") as f:
    data = json.load(f)


# ---- INIT PRESENTATION ----
prs = Presentation(TEMPLATE)

# Remove all existing template slides
sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst):
    rId = sldId.get(qn("r:id"))
    if rId:
        try:
            prs.part.drop_rel(rId)
        except KeyError:
            pass
    sldIdLst.remove(sldId)

# Use blank layout (index 37 in HPE template, fallback to last)
try:
    BLANK_LAYOUT = prs.slide_layouts[37]
except IndexError:
    BLANK_LAYOUT = prs.slide_layouts[-1]


# ---- HELPERS ----
def add_slide():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG_DARK
    return slide


def text_box(slide, left, top, width, height, text, *,
             font=FONT_BODY, size=14, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox


def multi_text(slide, left, top, width, height, runs_data, *, align=PP_ALIGN.LEFT):
    """runs_data: list of (text, font, size, color, bold)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for txt, fn, sz, clr, bld in runs_data:
        run = p.add_run()
        run.text = txt
        run.font.name = fn
        run.font.size = Pt(sz)
        run.font.color.rgb = clr
        run.font.bold = bld
    return txBox


def rounded_rect(slide, left, top, width, height, fill, *, border=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def accent_bar(slide, left, top, width, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.04),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ============================================================
# SLIDE 1 — COVER
# ============================================================
s = add_slide()
rounded_rect(s, 4.0, 1.0, 5.3, 0.45, RGBColor(0x00, 0xB8, 0x8C))
text_box(s, 4.0, 1.02, 5.3, 0.45, "STRATEGYZER PATTERN LIBRARY",
         font=FONT_HEADING, size=12, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)

text_box(s, 1.0, 2.0, 11.3, 0.8, "INVENT New Growth &", font=FONT_HEADING,
         size=36, color=HPE_GREEN, bold=True, align=PP_ALIGN.CENTER)
text_box(s, 1.0, 2.7, 11.3, 0.8, "SHIFT Existing Models", font=FONT_HEADING,
         size=36, color=HPE_PURPLE, bold=True, align=PP_ALIGN.CENTER)

text_box(s, 2.5, 3.8, 8.3, 1.0,
         "21 business model patterns from The Invincible Company — "
         "the canonical playbook for innovation portfolio strategy.",
         size=14, color=MUTED_TEXT, align=PP_ALIGN.CENTER)

# Stats row
stats = [("9", "INVENT Patterns", HPE_GREEN), ("27", "Flavors", HPE_BLUE),
         ("12", "SHIFT Patterns", HPE_PURPLE), ("21", "Total", WHITE)]
for i, (num, label, clr) in enumerate(stats):
    x = 2.0 + i * 2.6
    text_box(s, x, 5.0, 2.0, 0.6, num, font=FONT_HEADING, size=36,
             color=clr, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, x, 5.6, 2.0, 0.3, label, size=9, color=MUTED_TEXT,
             align=PP_ALIGN.CENTER)

text_box(s, 0.5, 6.8, 12.3, 0.4,
         "Source: Osterwalder, Pigneur, Etiemble, Smith — The Invincible Company (2020)",
         size=9, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2 — INVENT OVERVIEW
# ============================================================
s = add_slide()
text_box(s, 0.6, 0.3, 5, 0.3, "DIRECTION: INVENT", font=FONT_HEADING,
         size=10, color=HPE_GREEN, bold=True)
text_box(s, 0.6, 0.6, 10, 0.5, "9 Patterns for Creating New-Growth Business Models",
         font=FONT_HEADING, size=22, color=WHITE, bold=True)

invent_patterns = data["invent"]["patterns"]
cols = 3
row_h = 0.6
col_w = 3.9
start_y = 1.4
for idx, pat in enumerate(invent_patterns):
    c = idx % cols
    r = idx // cols
    x = 0.6 + c * (col_w + 0.15)
    y = start_y + r * (row_h + 0.55)
    cat_color = INVENT_CAT_COLORS.get(pat["category"], HPE_GREEN)
    rounded_rect(s, x, y, col_w, row_h + 0.4, CARD_BG, border=MID_GRAY)
    accent_bar(s, x, y, col_w, cat_color)
    text_box(s, x + 0.15, y + 0.08, col_w - 0.3, 0.3, pat["name"],
             font=FONT_HEADING, size=11, color=WHITE, bold=True)
    imp = pat.get("strategic_imperative", "")
    text_box(s, x + 0.15, y + 0.35, col_w - 0.3, 0.2, f"↑ {imp}",
             size=8, color=cat_color)
    flav_count = len(pat.get("flavors", []))
    text_box(s, x + 0.15, y + 0.6, col_w - 0.3, 0.2, f"{flav_count} flavors",
             size=8, color=MUTED_TEXT)


# ============================================================
# SLIDES 3–11 — ONE PER INVENT PATTERN
# ============================================================
for pat in invent_patterns:
    s = add_slide()
    cat_color = INVENT_CAT_COLORS.get(pat["category"], HPE_GREEN)

    # Header
    text_box(s, 0.6, 0.3, 5, 0.25, pat["category"].upper(),
             font=FONT_HEADING, size=10, color=cat_color, bold=True)
    text_box(s, 0.6, 0.55, 8, 0.5, pat["name"],
             font=FONT_HEADING, size=28, color=WHITE, bold=True)
    text_box(s, 0.6, 1.1, 5, 0.25, f"↑ {pat.get('strategic_imperative', '')}",
             size=12, color=cat_color, bold=True)

    # Description
    text_box(s, 0.6, 1.55, 7, 0.8, pat["description"],
             size=12, color=MUTED_TEXT)

    # Trigger question
    rounded_rect(s, 0.6, 2.5, 7, 0.7, RGBColor(0x2A, 0x3A, 0x35), border=RGBColor(0x00, 0x8F, 0x6E))
    text_box(s, 0.7, 2.45, 0.5, 0.2, "TRIGGER", size=7, color=HPE_GREEN, bold=True)
    text_box(s, 0.8, 2.65, 6.6, 0.5, pat.get("trigger_question", ""),
             size=11, color=HPE_CYAN)

    # Flavors grid
    flavors = pat.get("flavors", [])
    if flavors:
        text_box(s, 0.6, 3.4, 3, 0.25, f"FLAVORS ({len(flavors)})",
                 font=FONT_HEADING, size=10, color=HPE_BLUE, bold=True)

        flav_cols = min(len(flavors), 3)
        flav_w = (12.0 - 0.6) / flav_cols - 0.15
        for fi, flav in enumerate(flavors):
            fc = fi % flav_cols
            fr = fi // flav_cols
            fx = 0.6 + fc * (flav_w + 0.15)
            fy = 3.75 + fr * 1.6

            rounded_rect(s, fx, fy, flav_w, 1.45, CARD_BG, border=MID_GRAY)
            accent_bar(s, fx, fy, flav_w, HPE_BLUE)

            text_box(s, fx + 0.12, fy + 0.1, flav_w - 0.24, 0.25,
                     flav["name"], font=FONT_HEADING, size=10, color=HPE_BLUE, bold=True)
            text_box(s, fx + 0.12, fy + 0.35, flav_w - 0.24, 0.5,
                     flav["description"], size=8, color=MUTED_TEXT)

            examples = flav.get("examples", [])
            if examples:
                text_box(s, fx + 0.12, fy + 1.05, flav_w - 0.24, 0.25,
                         "e.g. " + ", ".join(examples), size=7, color=MID_GRAY)

    # Footer
    text_box(s, 0.5, 6.9, 12.3, 0.3,
             f"INVENT  ·  {pat['category']}  ·  {pat['name']}",
             size=8, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 12 — SHIFT OVERVIEW
# ============================================================
s = add_slide()
text_box(s, 0.6, 0.3, 5, 0.3, "DIRECTION: SHIFT", font=FONT_HEADING,
         size=10, color=HPE_PURPLE, bold=True)
text_box(s, 0.6, 0.6, 10, 0.5, "12 Patterns for Transforming Existing Business Models",
         font=FONT_HEADING, size=22, color=WHITE, bold=True)

shift_patterns = data["shift"]["patterns"]
cols = 3
for idx, pat in enumerate(shift_patterns):
    c = idx % cols
    r = idx // cols
    x = 0.6 + c * (col_w + 0.15)
    y = start_y + r * (row_h + 0.55)
    cat_color = SHIFT_CAT_COLORS.get(pat["category"], HPE_PURPLE)
    rounded_rect(s, x, y, col_w, row_h + 0.25, CARD_BG, border=MID_GRAY)
    accent_bar(s, x, y, col_w, cat_color)
    text_box(s, x + 0.15, y + 0.08, col_w - 0.3, 0.3, pat["name"],
             font=FONT_HEADING, size=10, color=WHITE, bold=True)
    ex = pat.get("example", "")
    text_box(s, x + 0.15, y + 0.4, col_w - 0.3, 0.2, f"📌 {ex}",
             size=8, color=cat_color)


# ============================================================
# SLIDES 13–24 — ONE PER SHIFT PATTERN
# ============================================================
for pat in shift_patterns:
    s = add_slide()
    cat_color = SHIFT_CAT_COLORS.get(pat["category"], HPE_PURPLE)

    # Header
    text_box(s, 0.6, 0.3, 5, 0.25, pat["category"].upper(),
             font=FONT_HEADING, size=10, color=cat_color, bold=True)
    text_box(s, 0.6, 0.55, 10, 0.5, pat["name"],
             font=FONT_HEADING, size=28, color=WHITE, bold=True)
    text_box(s, 0.6, 1.05, 1.5, 0.25, "↔ SHIFT",
             size=12, color=HPE_PURPLE, bold=True)

    # Description
    text_box(s, 0.6, 1.45, 8, 0.9, pat["description"],
             size=13, color=MUTED_TEXT)

    # Strategic reflection (green box)
    rounded_rect(s, 0.6, 2.6, 5.5, 1.0, RGBColor(0x2A, 0x3A, 0x35),
                 border=RGBColor(0x00, 0x8F, 0x6E))
    text_box(s, 0.75, 2.55, 3, 0.2, "STRATEGIC REFLECTION",
             size=7, color=HPE_GREEN, bold=True)
    text_box(s, 0.75, 2.8, 5.2, 0.7, pat.get("strategic_reflection", ""),
             size=11, color=HPE_CYAN)

    # Reverse reflection (orange box)
    rounded_rect(s, 6.4, 2.6, 5.5, 1.0, RGBColor(0x3A, 0x2D, 0x2A),
                 border=RGBColor(0x99, 0x55, 0x33))
    text_box(s, 6.55, 2.55, 3, 0.2, "REVERSE REFLECTION",
             size=7, color=HPE_ORANGE, bold=True)
    text_box(s, 6.55, 2.8, 5.2, 0.7, pat.get("reverse_reflection", ""),
             size=11, color=HPE_ORANGE)

    # Example badge
    ex = pat.get("example", "")
    if ex:
        rounded_rect(s, 0.6, 4.0, 2.5, 0.4, RGBColor(0x33, 0x30, 0x45),
                     border=RGBColor(0x7A, 0x6A, 0xCC))
        text_box(s, 0.7, 4.02, 2.3, 0.35, f"📌  {ex}",
                 size=10, color=HPE_PURPLE, bold=True, align=PP_ALIGN.CENTER)

    # Footer
    text_box(s, 0.5, 6.9, 12.3, 0.3,
             f"SHIFT  ·  {pat['category']}  ·  {pat['name']}",
             size=8, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# SAVE
# ============================================================
prs.save(OUTPUT)
print(f"✅ Saved {OUTPUT} — {len(prs.slides)} slides")
