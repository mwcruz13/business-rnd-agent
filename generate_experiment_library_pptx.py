"""
Generate Experiment Library PowerPoint from the HPE Dark template.

Slides:
  1.  Cover — 44 experiment cards overview
  2.  Evidence × Category Matrix (table)
  3.  Desirability — Weak Evidence cards
  4.  Desirability — Medium Evidence cards
  5.  Desirability — Strong Evidence cards
  6.  Feasibility — Weak Evidence cards
  7.  Feasibility — Medium Evidence cards
  8.  Feasibility — Strong Evidence cards
  9.  Viability — Weak Evidence cards
  10. Viability — Medium Evidence cards
  11. Viability — Strong Evidence cards
  12. Decision Tree
  13. Sequencing Rules
  14. Suggested Sequence Paths

Usage:
  cat generate_experiment_library_pptx.py | docker compose exec -T bmi-backend python -
"""
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# --- HPE Dark Theme Colors ---
BG_DARK     = RGBColor(0x29, 0x2D, 0x3A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
MUTED_TEXT  = RGBColor(0xB1, 0xB9, 0xBE)
MID_GRAY    = RGBColor(0x53, 0x5C, 0x66)
HPE_GREEN   = RGBColor(0x00, 0xE0, 0xAF)
HPE_PURPLE  = RGBColor(0x9B, 0x84, 0xFC)
HPE_BLUE    = RGBColor(0x65, 0xAE, 0xF9)
HPE_CYAN    = RGBColor(0x62, 0xE5, 0xF6)
HPE_ORANGE  = RGBColor(0xF2, 0x6B, 0x43)
HPE_YELLOW  = RGBColor(0xF2, 0xC9, 0x4C)
CARD_BG     = RGBColor(0x33, 0x3A, 0x4A)

FONT_HEADING = "HPE Graphik Semibold"
FONT_BODY    = "HPE Graphik"

# Category → accent color
CAT_COLORS = {
    "Desirability": HPE_GREEN,
    "Feasibility": HPE_BLUE,
    "Viability": HPE_PURPLE,
}

EVIDENCE_COLORS = {
    "Weak": HPE_BLUE,
    "Medium": HPE_YELLOW,
    "Strong": HPE_GREEN,
}

# --- Paths ---
TEMPLATE   = "hpe_dark_template.pptx"
DATA_FILE  = "backend/app/patterns/experiment-library.json"
OUTPUT     = "reports/Experiment_Library.pptx"

# ---- LOAD DATA ----
with open(DATA_FILE, encoding="utf-8") as f:
    data = json.load(f)

# ---- INIT PRESENTATION ----
prs = Presentation(TEMPLATE)

# Remove all existing template slides
sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst):
    rId = sldId.get(qn("r:id"))
    if rId:
        try:
            prs.part.drop_rel(rId)
        except KeyError:
            pass
    sldIdLst.remove(sldId)

try:
    BLANK_LAYOUT = prs.slide_layouts[37]
except IndexError:
    BLANK_LAYOUT = prs.slide_layouts[-1]


# ---- HELPERS ----
def add_slide():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG_DARK
    return slide


def text_box(slide, left, top, width, height, text, *,
             font=FONT_BODY, size=14, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox


def multi_line_box(slide, left, top, width, height, lines, *, align=PP_ALIGN.LEFT):
    """lines: list of (text, font, size, color, bold)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, fn, sz, clr, bld) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = txt
        run.font.name = fn
        run.font.size = Pt(sz)
        run.font.color.rgb = clr
        run.font.bold = bld
        p.space_after = Pt(4)
    return txBox


def rounded_rect(slide, left, top, width, height, fill, *, border=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def accent_bar(slide, left, top, width, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.04),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ============================================================
# SLIDE 1 — COVER
# ============================================================
experiments = data["experiments"]

s = add_slide()
rounded_rect(s, 4.0, 0.8, 5.3, 0.45, RGBColor(0x00, 0xB8, 0x8C))
text_box(s, 4.0, 0.82, 5.3, 0.45, "TESTING BUSINESS IDEAS",
         font=FONT_HEADING, size=12, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)

text_box(s, 1.0, 1.8, 11.3, 0.8, "44 Experiment Cards for",
         font=FONT_HEADING, size=36, color=HPE_GREEN, bold=True, align=PP_ALIGN.CENTER)
text_box(s, 1.0, 2.5, 11.3, 0.8, "De-Risking Business Ideas",
         font=FONT_HEADING, size=36, color=HPE_PURPLE, bold=True, align=PP_ALIGN.CENTER)

text_box(s, 2.5, 3.6, 8.3, 1.0,
         "The canonical experiment card library from Bland & Osterwalder — "
         "organized by Desirability, Feasibility, and Viability with "
         "evidence strength progression from Weak → Medium → Strong.",
         size=13, color=MUTED_TEXT, align=PP_ALIGN.CENTER)

# Count by evidence
weak_count = sum(1 for e in experiments if e["evidence_strength"] == "Weak")
med_count = sum(1 for e in experiments if e["evidence_strength"] == "Medium")
strong_count = sum(1 for e in experiments if e["evidence_strength"] == "Strong")

stats = [
    ("44", "Experiment Cards", WHITE),
    (str(weak_count), "Weak Evidence", HPE_BLUE),
    (str(med_count), "Medium Evidence", HPE_YELLOW),
    (str(strong_count), "Strong Evidence", HPE_GREEN),
]
for i, (num, label, clr) in enumerate(stats):
    x = 2.0 + i * 2.6
    text_box(s, x, 5.0, 2.0, 0.6, num, font=FONT_HEADING, size=36,
             color=clr, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, x, 5.6, 2.0, 0.3, label, size=9, color=MUTED_TEXT,
             align=PP_ALIGN.CENTER)

text_box(s, 0.5, 6.8, 12.3, 0.4,
         "Source: David J. Bland & Alexander Osterwalder — Testing Business Ideas (2019)",
         size=9, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2 — EVIDENCE × CATEGORY MATRIX
# ============================================================
s = add_slide()
text_box(s, 0.6, 0.3, 8, 0.3, "EVIDENCE × CATEGORY MATRIX",
         font=FONT_HEADING, size=10, color=HPE_ORANGE, bold=True)
text_box(s, 0.6, 0.6, 10, 0.5, "Cards by Evidence Strength and DVF Category",
         font=FONT_HEADING, size=22, color=WHITE, bold=True)

matrix = data["matrix_by_evidence_and_category"]
# Table header
col_w = [1.8, 3.3, 3.3, 3.3]
col_x = [0.6]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w + 0.1)

y = 1.4
# Header row
for i, (header, clr) in enumerate([
    ("Evidence", WHITE), ("Desirability", HPE_GREEN),
    ("Feasibility", HPE_BLUE), ("Viability", HPE_PURPLE)
]):
    rounded_rect(s, col_x[i], y, col_w[i], 0.4, MID_GRAY)
    text_box(s, col_x[i], y + 0.02, col_w[i], 0.36, header,
             font=FONT_HEADING, size=10, color=clr, bold=True, align=PP_ALIGN.CENTER)

# Data rows
for level in ["Weak", "Medium", "Strong"]:
    y += 0.55
    ev_clr = EVIDENCE_COLORS[level]
    rounded_rect(s, col_x[0], y, col_w[0], 0.45, CARD_BG)
    text_box(s, col_x[0], y + 0.04, col_w[0], 0.4, level,
             font=FONT_HEADING, size=11, color=ev_clr, bold=True, align=PP_ALIGN.CENTER)
    for j, cat in enumerate(["Desirability", "Feasibility", "Viability"]):
        cards = matrix[level].get(cat, [])
        cell_text = ", ".join(cards) if cards else "—"
        rounded_rect(s, col_x[j+1], y, col_w[j+1], 0.45, CARD_BG)
        text_box(s, col_x[j+1] + 0.08, y + 0.04, col_w[j+1] - 0.16, 0.4,
                 cell_text, size=8, color=MUTED_TEXT)

text_box(s, 0.6, 6.8, 12.3, 0.4,
         "Weak → Medium → Strong: progress experiments as evidence accumulates",
         size=9, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDES 3–11 — CARDS BY CATEGORY + EVIDENCE
# ============================================================
for cat in ["Desirability", "Feasibility", "Viability"]:
    cat_color = CAT_COLORS[cat]
    for strength in ["Weak", "Medium", "Strong"]:
        ev_color = EVIDENCE_COLORS[strength]
        cards = [e for e in experiments
                 if e["category"] == cat and e["evidence_strength"] == strength]
        if not cards:
            continue

        s = add_slide()
        text_box(s, 0.6, 0.3, 5, 0.3, f"{cat.upper()} — {strength.upper()} EVIDENCE",
                 font=FONT_HEADING, size=10, color=cat_color, bold=True)
        text_box(s, 0.6, 0.6, 10, 0.5, f"{len(cards)} Experiment Cards",
                 font=FONT_HEADING, size=22, color=WHITE, bold=True)

        # Card grid — 2 columns
        col_count = 2
        card_w = 5.9
        card_h = 1.6
        gap_x = 0.5
        gap_y = 0.2
        start_y = 1.3
        start_x = 0.6

        for idx, exp in enumerate(cards):
            col = idx % col_count
            row = idx // col_count
            cx = start_x + col * (card_w + gap_x)
            cy = start_y + row * (card_h + gap_y)

            # Overflow to new slide
            if cy + card_h > 7.2:
                s = add_slide()
                text_box(s, 0.6, 0.3, 5, 0.3,
                         f"{cat.upper()} — {strength.upper()} EVIDENCE (cont.)",
                         font=FONT_HEADING, size=10, color=cat_color, bold=True)
                row = 0
                cy = 0.7
                col = idx % col_count
                cx = start_x + col * (card_w + gap_x)

            rounded_rect(s, cx, cy, card_w, card_h, CARD_BG, border=MID_GRAY)
            # Accent bar at top
            shape = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(cx), Inches(cy), Inches(card_w), Inches(0.04),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = cat_color
            shape.line.fill.background()

            # Card ID + Name
            text_box(s, cx + 0.12, cy + 0.1, 0.5, 0.2,
                     f"#{exp['id']}", size=8, color=MID_GRAY, bold=True)
            text_box(s, cx + 0.5, cy + 0.1, card_w - 0.7, 0.25,
                     exp["name"], font=FONT_HEADING, size=11, color=WHITE, bold=True)

            # What it tests
            text_box(s, cx + 0.12, cy + 0.4, card_w - 0.24, 0.15,
                     "WHAT IT TESTS", size=7, color=MID_GRAY, bold=True)
            text_box(s, cx + 0.12, cy + 0.55, card_w - 0.24, 0.4,
                     exp["what_it_tests"], size=8, color=MUTED_TEXT)

            # Best used when
            text_box(s, cx + 0.12, cy + 0.95, card_w - 0.24, 0.15,
                     "BEST USED WHEN", size=7, color=MID_GRAY, bold=True)
            text_box(s, cx + 0.12, cy + 1.1, card_w - 0.24, 0.45,
                     exp["best_used_when"], size=8, color=MUTED_TEXT)


# ============================================================
# SLIDE 12 — DECISION TREE
# ============================================================
tree = data["decision_tree"]
s = add_slide()
text_box(s, 0.6, 0.3, 5, 0.3, "EXPERIMENT DECISION TREE",
         font=FONT_HEADING, size=10, color=HPE_CYAN, bold=True)
text_box(s, 0.6, 0.6, 10, 0.5, "Start with the question, then choose the context",
         font=FONT_HEADING, size=18, color=WHITE, bold=True)

branches = tree["branches"]
y = 1.4
for branch in branches:
    bcat = branch["category"]
    bcolor = CAT_COLORS.get(bcat, HPE_CYAN)

    # Branch question
    rounded_rect(s, 0.6, y, 12.1, 0.35, CARD_BG)
    shape = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(y), Inches(0.06), Inches(0.35),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bcolor
    shape.line.fill.background()
    text_box(s, 0.8, y + 0.02, 11.5, 0.3, branch["question"],
             font=FONT_HEADING, size=11, color=WHITE, bold=True)
    y += 0.42

    # Sub branches
    for sb in branch["sub_branches"]:
        if y > 6.8:
            s = add_slide()
            text_box(s, 0.6, 0.3, 5, 0.3, "DECISION TREE (cont.)",
                     font=FONT_HEADING, size=10, color=HPE_CYAN, bold=True)
            y = 0.7

        text_box(s, 1.0, y, 3.5, 0.25, sb["context"],
                 size=9, color=HPE_CYAN, bold=True)
        cards_text = " → ".join(sb["cards"])
        text_box(s, 4.6, y, 8.0, 0.25, cards_text,
                 size=9, color=MUTED_TEXT)
        y += 0.3

    y += 0.15


# ============================================================
# SLIDE 13 — SEQUENCING RULES
# ============================================================
s = add_slide()
text_box(s, 0.6, 0.3, 5, 0.3, "SEQUENCING RULES",
         font=FONT_HEADING, size=10, color=HPE_GREEN, bold=True)
text_box(s, 0.6, 0.6, 10, 0.5, "Principles for Ordering Experiments",
         font=FONT_HEADING, size=22, color=WHITE, bold=True)

rules = data["sequencing_rules"]
y = 1.4
for i, rule in enumerate(rules):
    rounded_rect(s, 0.6, y, 12.1, 0.5, CARD_BG)
    text_box(s, 0.8, y + 0.06, 0.4, 0.4, str(i + 1),
             font=FONT_HEADING, size=18, color=HPE_GREEN, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, 1.4, y + 0.08, 11.0, 0.38, rule,
             size=11, color=MUTED_TEXT)
    y += 0.58


# ============================================================
# SLIDE 14 — SUGGESTED SEQUENCE PATHS
# ============================================================
paths = data["suggested_sequence_paths"]
s = add_slide()
text_box(s, 0.6, 0.3, 5, 0.3, "SUGGESTED SEQUENCE PATHS",
         font=FONT_HEADING, size=10, color=HPE_PURPLE, bold=True)
text_box(s, 0.6, 0.6, 10, 0.5, "Recommended Experiment Progressions by Category",
         font=FONT_HEADING, size=22, color=WHITE, bold=True)

y = 1.4
for cat_key, routes in paths.items():
    cat_color = CAT_COLORS.get(cat_key.capitalize(),
                                CAT_COLORS.get(cat_key.title(), HPE_GREEN))
    text_box(s, 0.6, y, 3.0, 0.3, cat_key.upper(),
             font=FONT_HEADING, size=12, color=cat_color, bold=True)
    y += 0.35

    for route_name, cards in routes.items():
        if y > 6.5:
            s = add_slide()
            text_box(s, 0.6, 0.3, 5, 0.3, "SEQUENCE PATHS (cont.)",
                     font=FONT_HEADING, size=10, color=HPE_PURPLE, bold=True)
            y = 0.7

        label = route_name.replace("_", " ").title()
        text_box(s, 0.8, y, 2.8, 0.25, label,
                 size=9, color=MID_GRAY, bold=True)

        arrow_text = " → ".join(cards)
        text_box(s, 3.8, y, 8.8, 0.25, arrow_text,
                 size=9, color=MUTED_TEXT)
        y += 0.3

    y += 0.2


# ---- SAVE ----
Path("reports").mkdir(exist_ok=True)
prs.save(OUTPUT)
print(f"Saved {OUTPUT} with {len(prs.slides)} slides")
